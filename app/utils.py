from app import slack_client as slack_client_module
from app import orquesta_client as orquesta_client_module
from io import BytesIO
import threading
import shlex
import logging
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
import requests
from slack_sdk.errors import SlackApiError
import os

def handle_app_mention(event):
    logging.info(f"Handling event: {event}")  # Log the event being handled
    # Ignore events where the user is the bot itself
    bot_user_id = slack_client_module.bot_user_id
    if event.get('user') == bot_user_id:
        return
    
    if event.get('type') == 'app_mention':
        prompt_user = event.get('text', '').split('>')[1].strip()
    elif event.get('type') == 'message' and event.get('channel_type') == 'im':
        prompt_user = event.get('text', '').strip()

    files = event.get('files', [])
    text_content = None  # Initialize text_content

    if files:
        for file_info in files:
            text_content = handle_file(file_info, event)  # Capture the returned text_content
            if text_content:  # If text_content is available, break the loop
                break

    # Use slack_client from the slack_client module
    slack_client_module.slack_client.chat_postMessage(
        channel=event['channel'],
        thread_ts=event['ts'],
        text="Processing your request, please wait..."
    )

    threading.Thread(target=query_orquesta, args=(event, prompt_user, text_content)).start()

def handle_file(file_info, event):
    file_id = file_info.get('id')
    text_content = None  # Initialize text_content
    try:
        file_url = file_info.get('url_private_download')  # Use the direct download URL
        file_content = download_file(file_url)
        text_content = process_file_content(file_content, event)  # Get text_content from the file
    except SlackApiError as e:
        logging.error(f"Error getting file info: {e}")
    return text_content  # Return the extracted text content
    
def download_file(file_url):
    headers = {'Authorization': f'Bearer {os.getenv("SLACK_BOT_TOKEN")}'}
    response = requests.get(file_url, headers=headers)
    if response.status_code == 200:
        logging.info(f"File downloaded successfully: {file_url}")  # Log successful download
        return response.content
    else:
        logging.error(f"Error downloading file: {response.status_code}, {response.text}")  # Log download error
        return None

def process_file_content(file_content, event):
    file_info = event.get('files', [])[0]  # Assuming there's at least one file
    file_type = file_info.get('filetype')

    text_content = None
    if file_type == 'pdf':
        text_content = extract_text_from_pdf(file_content)
    elif file_type in ['doc', 'docx']:
        text_content = extract_text_from_docx(file_content)
    elif file_type in ['ppt', 'pptx']:
        text_content = extract_text_from_pptx(file_content)

    if text_content:
        logging.info(f"Extracted text content: {text_content[:10]}")  # Log the first 10 characters
    else:
        logging.error(f"Could not extract text from the file of type {file_type}.")
    return text_content

def extract_text_from_pdf(file_content):
    try:
        reader = PdfReader('temp_file.pdf')
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        logging.error(f"Error extracting text from PDF: {e}")
        return None

def extract_text_from_docx(file_content):
    try:
        # Convert bytes content to a file-like object
        file_stream = BytesIO(file_content)
        doc = Document(file_stream)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    except Exception as e:
        logging.error(f"Error extracting text from DOCX: {e}")
        return None

def extract_text_from_pptx(file_content):
    try:
        # Convert bytes content to a file-like object
        file_stream = BytesIO(file_content)
        ppt = Presentation(file_stream)
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logging.error(f"Error extracting text from PPTX: {e}")
        return None

def parse_command_arguments(command_text):
    try:
        return shlex.split(command_text)
    except ValueError as e:
        raise ValueError(f"Error parsing arguments: {e}. Make sure to enclose each argument with double quotes.")

def post_error_message(channel_id, ts, message):
    if slack_client_module.slack_client is None:
        logging.error("Slack client has not been initialized.")
        return
    slack_client_module.slack_client.chat_postMessage(
        channel=channel_id,
        thread_ts=ts,
        text=message
    )

def query_orquesta(event, prompt_user, text_content):
    logging.info(f"Invoking Orquesta with event: {event}, prompt_user: {prompt_user}, text_content: {text_content}")  # Log the Orquesta invocation
    # Check if text_content is empty
    if not text_content:
        # Invoke the Orquesta deployment
        deployment = orquesta_client_module.client.deployments.invoke(
            key="pierre-fabre-slack-app",
            context={
                "doc": False
            },
            inputs={
                "prompt": prompt_user
            }
        )
    else:
        # Invoke the Orquesta deployment
        deployment = orquesta_client_module.client.deployments.invoke(
            key="pierre-fabre-slack-app",
            context={
                "environments": [],
                "doc": None
            },
            inputs={
                "doc": text_content,
                "prompt": prompt_user
            }
        )

    logging.info(f"Posting message to thread: {deployment.choices[0].message.content}")  # Log the message being posted
    slack_client_module.slack_client.chat_postMessage(
        channel=event['channel'],
        thread_ts=event['ts'],  # Ensure this is the original message timestamp
        text=deployment.choices[0].message.content
    )
    